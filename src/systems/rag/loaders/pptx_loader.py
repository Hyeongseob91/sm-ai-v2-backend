"""PowerPoint 파일 로더 - python-pptx 기반"""

from typing import List
from langchain_core.documents import Document
from pptx import Presentation


class PowerPointLoader:
    """
    python-pptx를 사용한 PowerPoint 파일 로더.
    각 슬라이드를 별도의 Document로 변환합니다.
    """

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        """PowerPoint 파일을 로드하여 Document 리스트로 반환합니다."""
        documents = []

        prs = Presentation(self.file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            content = self._extract_slide_text(slide)

            if content.strip():
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": self.file_path,
                        "slide_number": slide_num,
                        "file_type": "pptx",
                    },
                )
                documents.append(doc)

        return documents

    def _extract_slide_text(self, slide) -> str:
        """슬라이드에서 모든 텍스트를 추출합니다."""
        text_parts = []

        for shape in slide.shapes:
            # 텍스트 프레임이 있는 도형에서 텍스트 추출
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text.strip())

            # 테이블에서 텍스트 추출
            if shape.has_table:
                table_text = self._extract_table_text(shape.table)
                if table_text:
                    text_parts.append(table_text)

        return "\n".join(text_parts)

    def _extract_table_text(self, table) -> str:
        """테이블에서 텍스트를 추출합니다."""
        rows = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_cells.append(cell_text)
            rows.append(" | ".join(row_cells))
        return "\n".join(rows)
